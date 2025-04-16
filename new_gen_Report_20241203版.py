import PyPDF2, cv2, os, pptx
from pptx.util import Inches

# Load the PowerPoint template
prs = pptx.Presentation('MXPT_report_基礎型_20241203版.pptx')

# Prompt for PDF filenames
print("依以下順序輸入PDF檔名:\n1.上肢PDF\n2.下肢PDF\n3.張眼平衡PDF\n4.閉眼平衡PDF")

# Process the upper limb PDF
fileName1 = input("\n >> 上肢PDF名稱: ")
pdf_file = open(fileName1, 'rb')
pdf_reader = PyPDF2.PdfReader(pdf_file)

# Define image processing parameters
pic_name = ["kf_upper_1.jpg", "kf_upper_3.jpg", "kf_upper_2.jpg"]
page_name = ["1.jpg", "2.jpg", "3.jpg"]
page = [0, 0, 1]
crop_size = [[1166, 2500, 80, 2900], [2734, 3700, 80, 2900], [132, 2049, 80, 2900]]

# Extract and crop images from the PDF
i = 0
for page_num in page:
    page_obj = pdf_reader.pages[page_num]
    page_objs = page_obj['/Resources']['/XObject'].get_object()
    for obj_name in page_objs:
        img_data = page_objs[obj_name].get_data()
        with open(page_name[i], 'wb') as img_file:
            img_file.write(img_data)
        img = cv2.imread(page_name[i])
        cropped_img = img[int(crop_size[i][0]):int(crop_size[i][1]), int(crop_size[i][2]):int(crop_size[i][3])]
        cv2.imwrite(pic_name[i], cropped_img)
    i += 1

# Clean up intermediate files
for jpg_name in page_name:
    os.remove(jpg_name)

# Load the new images into the PowerPoint slide
new_pptx_img = [
    pptx.parts.image.Image.from_file("kf_upper_3.jpg"),
    pptx.parts.image.Image.from_file("kf_upper_2.jpg"),
    pptx.parts.image.Image.from_file("kf_upper_1.jpg"),
    pptx.parts.image.Image.from_file("vd_upper_3.jpg"),
    pptx.parts.image.Image.from_file("vd_upper_2.jpg"),
    pptx.parts.image.Image.from_file("vd_upper_1.jpg"),
]
target_texts = ["kf_upper_3.jpg", "kf_upper_2.jpg", "kf_upper_1.jpg", "vd_upper_3.jpg", "vd_upper_2.jpg", "vd_upper_1.jpg"]
ppt_page = 1

# Replace images in shapes based on matching text
for i, new_image in enumerate(new_pptx_img):
    try:
        # Iterate over shapes in the target slide
        for shape in prs.slides[ppt_page].shapes:
            # Check if the shape contains the target text
            if shape.has_text_frame and shape.text.strip() == target_texts[i]:
                # Get the position and size of the matching shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save the new image as a temporary file
                temp_image_path = f"temp_image_{i}.jpg"
                with open(temp_image_path, "wb") as temp_file:
                    temp_file.write(new_image._blob)

                # Add a new picture shape at the same position and size
                picture = prs.slides[ppt_page].shapes.add_picture(temp_image_path, left, top, width, height)

                # Remove the original shape
                prs.slides[ppt_page].shapes._spTree.remove(shape._element)

                # Delete the temporary image file
                os.remove(temp_image_path)

                # Stop searching once the shape is replaced
                break

    except Exception as err:
        print(f"Error updating shape {i} on slide {ppt_page}: {err}")
        pass

# Save the updated PowerPoint
prs.save('MXPT_report_基礎型_20241203版.pptx')



fileName2 = input("\n >> 下肢PDF名稱: ")
pdf_file = open(fileName2, 'rb')
pdf_reader = PyPDF2.PdfReader(pdf_file)
upOrDown = ""
pic_name = ["kf_lower_1.jpg", "kf_lower_3.jpg", "kf_lower_2.jpg"]
upOrDown = "d"
page_name = ["1.jpg", "2.jpg", "3.jpg"]
page = [0,0,1]
crop_size = [ [1100, 2400, 80, 2900], [3550, 4100, 80, 2900], [1400, 2900, 90, 2910] ]
# Extract and crop images from the PDF
i = 0
for page_num in page:
    page_obj = pdf_reader.pages[page_num]
    page_objs = page_obj['/Resources']['/XObject'].get_object()
    for obj_name in page_objs:
        img_data = page_objs[obj_name].get_data()
        with open(page_name[i], 'wb') as img_file:
            img_file.write(img_data)
        img = cv2.imread(page_name[i])
        cropped_img = img[int(crop_size[i][0]):int(crop_size[i][1]), int(crop_size[i][2]):int(crop_size[i][3])]
        cv2.imwrite(pic_name[i], cropped_img)
    i += 1

# Clean up intermediate files
for jpg_name in page_name:
    os.remove(jpg_name)

# Load the new images into the PowerPoint slide
new_pptx_img = [
    pptx.parts.image.Image.from_file("kf_lower_3.jpg"),
    pptx.parts.image.Image.from_file("kf_lower_2.jpg"),
    pptx.parts.image.Image.from_file("kf_lower_1.jpg"),
    pptx.parts.image.Image.from_file("vd_lower_3.jpg"),
    pptx.parts.image.Image.from_file("vd_lower_2.jpg"),
    pptx.parts.image.Image.from_file("vd_lower_1.jpg"),
]
target_texts = ["kf_lower_3.jpg", "kf_lower_2.jpg", "kf_lower_1.jpg", "vd_lower_3.jpg", "vd_lower_2.jpg", "vd_lower_1.jpg"]
ppt_page = 2

# Replace images in shapes based on matching text
for i, new_image in enumerate(new_pptx_img):
    try:
        # Iterate over shapes in the target slide
        for shape in prs.slides[ppt_page].shapes:
            # Check if the shape contains the target text
            if shape.has_text_frame and shape.text.strip() == target_texts[i]:
                # Get the position and size of the matching shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save the new image as a temporary file
                temp_image_path = f"temp_image_{i}.jpg"
                with open(temp_image_path, "wb") as temp_file:
                    temp_file.write(new_image._blob)

                # Add a new picture shape at the same position and size
                picture = prs.slides[ppt_page].shapes.add_picture(temp_image_path, left, top, width, height)

                # Remove the original shape
                prs.slides[ppt_page].shapes._spTree.remove(shape._element)

                # Delete the temporary image file
                os.remove(temp_image_path)

                # Stop searching once the shape is replaced
                break

    except Exception as err:
        print(f"Error updating shape {i} on slide {ppt_page}: {err}")
        pass

# Save the updated PowerPoint
prs.save('MXPT_report_基礎型_20241203版.pptx')



#stiffness
def find_stiffness_graph(directory="."):
    """
    Search for a file in the given directory containing the keyword 'stiffness'.
    Returns the full path of the first matching file or None if no file is found.
    """
    for file in os.listdir(directory):
        if "stiffness" in file.lower() and file.endswith((".png", ".jpg")):  # Check for valid image files
            return os.path.join(directory, file)
    return None

# Find the stiffness graph file
stiffness_graph_path = find_stiffness_graph()

if stiffness_graph_path:
    print(f"Found stiffness graph: {stiffness_graph_path}")

    # Replace the stiffness graph on slide 2
    ppt_page = 2  # Assuming slide 2 contains the "stiffness" graph
    for shape in prs.slides[ppt_page].shapes:
        try:
            # Check if the shape contains the target "stiffness" text
            if shape.has_text_frame and "stiffness" in shape.text.lower():
                # Get the position and size of the matching shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Add the stiffness graph at the same position and size
                prs.slides[ppt_page].shapes.add_picture(stiffness_graph_path, left, top, width, height)

                # Remove the original placeholder shape
                prs.slides[ppt_page].shapes._spTree.remove(shape._element)

                print("Stiffness graph replaced successfully.")
                break

        except Exception as err:
            print(f"Error replacing stiffness graph on slide {ppt_page}: {err}")
            pass

    # Save the updated PowerPoint
    prs.save('MXPT_report_基礎型_20241203版.pptx')
else:
    print("No stiffness graph found in the current directory.")





fileName3 = input("\n >> 張眼平衡PDF名稱: ")
pdf_file = open(fileName3, 'rb')
pdf_reader = PyPDF2.PdfReader(pdf_file)
upOrDown = ""
pic_name = ["kf_balance_open_1.jpg", "kf_balance_open_2.jpg"]
upOrDown = "b"
page_name = ["1.jpg", "2.jpg"]
page = [0,0]
crop_size = [ [803, 2148, 80, 2900], [2138, 3132, 550, 2430] ]
# Extract and crop images from the PDF
i = 0
for page_num in page:
    page_obj = pdf_reader.pages[page_num]
    page_objs = page_obj['/Resources']['/XObject'].get_object()
    for obj_name in page_objs:
        img_data = page_objs[obj_name].get_data()
        with open(page_name[i], 'wb') as img_file:
            img_file.write(img_data)
        img = cv2.imread(page_name[i])
        cropped_img = img[int(crop_size[i][0]):int(crop_size[i][1]), int(crop_size[i][2]):int(crop_size[i][3])]
        cv2.imwrite(pic_name[i], cropped_img)
    i += 1

# Clean up intermediate files
for jpg_name in page_name:
    os.remove(jpg_name)

# Load the new images into the PowerPoint slide
new_pptx_img = [

    pptx.parts.image.Image.from_file("kf_balance_open_2.jpg"),
    pptx.parts.image.Image.from_file("kf_balance_open_1.jpg"),
    pptx.parts.image.Image.from_file("vd_balance_open_2.jpg"),
    pptx.parts.image.Image.from_file("vd_balance_open_1.jpg"),
]
target_texts = [ "kf_balance_open_2.jpg", "kf_balance_open_1.jpg", "vd_balance_open_2.jpg", "vd_balance_open_1.jpg"]
ppt_page = 3

# Replace images in shapes based on matching text
for i, new_image in enumerate(new_pptx_img):
    try:
        # Iterate over shapes in the target slide
        for shape in prs.slides[ppt_page].shapes:
            # Check if the shape contains the target text
            if shape.has_text_frame and shape.text.strip() == target_texts[i]:
                # Get the position and size of the matching shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save the new image as a temporary file
                temp_image_path = f"temp_image_{i}.jpg"
                with open(temp_image_path, "wb") as temp_file:
                    temp_file.write(new_image._blob)

                # Add a new picture shape at the same position and size
                picture = prs.slides[ppt_page].shapes.add_picture(temp_image_path, left, top, width, height)

                # Remove the original shape
                prs.slides[ppt_page].shapes._spTree.remove(shape._element)

                # Delete the temporary image file
                os.remove(temp_image_path)

                # Stop searching once the shape is replaced
                break

    except Exception as err:
        print(f"Error updating shape {i} on slide {ppt_page}: {err}")
        pass

# Save the updated PowerPoint
prs.save('MXPT_report_基礎型_20241203版.pptx')






fileName4 = input("\n >> 閉眼平衡PDF名稱: ")
pdf_file = open(fileName4, 'rb')
pdf_reader = PyPDF2.PdfReader(pdf_file)
upOrDown = ""
pic_name = ["kf_balance_close_1.jpg", "kf_balance_close_2.jpg"]
upOrDown = "b"
page_name = ["1.jpg", "2.jpg"]
page = [0,0]
crop_size = [ [803, 2148, 80, 2900], [2138, 3132, 550, 2430] ]
# Extract and crop images from the PDF
i = 0
for page_num in page:
    page_obj = pdf_reader.pages[page_num]
    page_objs = page_obj['/Resources']['/XObject'].get_object()
    for obj_name in page_objs:
        img_data = page_objs[obj_name].get_data()
        with open(page_name[i], 'wb') as img_file:
            img_file.write(img_data)
        img = cv2.imread(page_name[i])
        cropped_img = img[int(crop_size[i][0]):int(crop_size[i][1]), int(crop_size[i][2]):int(crop_size[i][3])]
        cv2.imwrite(pic_name[i], cropped_img)
    i += 1

# Clean up intermediate files
for jpg_name in page_name:
    os.remove(jpg_name)

# Load the new images into the PowerPoint slide
new_pptx_img = [

    pptx.parts.image.Image.from_file("kf_balance_close_2.jpg"),
    pptx.parts.image.Image.from_file("kf_balance_close_1.jpg"),
    pptx.parts.image.Image.from_file("vd_balance_close_2.jpg"),
    pptx.parts.image.Image.from_file("vd_balance_close_1.jpg"),
]
target_texts = [ "kf_balance_close_2.jpg", "kf_balance_close_1.jpg", "vd_balance_close_2.jpg", "vd_balance_close_1.jpg"]
ppt_page = 3

# Replace images in shapes based on matching text
for i, new_image in enumerate(new_pptx_img):
    try:
        # Iterate over shapes in the target slide
        for shape in prs.slides[ppt_page].shapes:
            # Check if the shape contains the target text
            if shape.has_text_frame and shape.text.strip() == target_texts[i]:
                # Get the position and size of the matching shape
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Save the new image as a temporary file
                temp_image_path = f"temp_image_{i}.jpg"
                with open(temp_image_path, "wb") as temp_file:
                    temp_file.write(new_image._blob)

                # Add a new picture shape at the same position and size
                picture = prs.slides[ppt_page].shapes.add_picture(temp_image_path, left, top, width, height)

                # Remove the original shape
                prs.slides[ppt_page].shapes._spTree.remove(shape._element)

                # Delete the temporary image file
                os.remove(temp_image_path)

                # Stop searching once the shape is replaced
                break

    except Exception as err:
        print(f"Error updating shape {i} on slide {ppt_page}: {err}")
        pass

# Save the updated PowerPoint
prs.save('MXPT_report_基礎型_20241203版.pptx')